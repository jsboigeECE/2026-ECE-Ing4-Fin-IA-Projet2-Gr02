"""
Génération des slides — Groupe 02 — Version simple et visuelle.
Usage : python slides/generate_slides.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ─────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1A, 0x23, 0x7E)
BLUE      = RGBColor(0x19, 0x76, 0xD2)
LIGHTBLUE = RGBColor(0xE3, 0xF2, 0xFD)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x37, 0x47, 0x4F)
LIGHTGRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN     = RGBColor(0x2E, 0x7D, 0x32)
RED       = RGBColor(0xC6, 0x28, 0x28)
ORANGE    = RGBColor(0xE6, 0x51, 0x00)
YELLOW    = RGBColor(0xF9, 0xA8, 0x25)

W = Inches(13.33)
H = Inches(7.5)

FIGURES = Path(__file__).parent.parent / "docs" / "figures"


# ── Helpers ─────────────────────────────────────────────────────────────────

def new_prs():
    p = Presentation()
    p.slide_width  = W
    p.slide_height = H
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def box(slide, x, y, w, h, fill=None, border=None, radius=False):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid() if fill else s.fill.background()
    if fill:
        s.fill.fore_color.rgb = fill
    s.line.fill.background() if not border else None
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=20, bold=False, color=GRAY,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def header(slide, title, subtitle=None):
    """Bandeau bleu en haut."""
    h_height = Inches(1.3) if not subtitle else Inches(1.55)
    box(slide, 0, 0, W, h_height, fill=NAVY)
    txt(slide, title,
        Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.85),
        size=34, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.55),
            size=16, color=RGBColor(0xBB, 0xDE, 0xFB))

def footer(slide):
    txt(slide, "ECE Paris · Groupe 02 · Bayesian Sports Analytics · 31 mars 2026",
        Inches(0.3), Inches(7.1), Inches(12.5), Inches(0.35),
        size=9, color=RGBColor(0xB0, 0xBE, 0xC5))

def stat_box(slide, number, label, x, y, w=Inches(2.8), h=Inches(1.6),
             bg=BLUE, num_color=WHITE, lbl_color=WHITE):
    """Grande boîte avec un chiffre clé."""
    box(slide, x, y, w, h, fill=bg)
    txt(slide, number, x, y + Inches(0.15), w, Inches(0.85),
        size=38, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    txt(slide, label, x, y + Inches(0.95), w, Inches(0.55),
        size=14, color=lbl_color, align=PP_ALIGN.CENTER)

def bullets(slide, items, x, y, w, h, base_size=20):
    """Liste de bullets propre."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        sub = item.startswith("  ")
        p.space_before = Pt(3 if sub else 10)
        r = p.add_run()
        r.text = ("      " if sub else "▶  ") + item.strip()
        r.font.size  = Pt((base_size - 3) if sub else base_size)
        r.font.bold  = not sub
        r.font.color.rgb = GRAY if sub else NAVY

def img(slide, path, x, y, w, h):
    p = Path(path)
    if p.exists():
        slide.shapes.add_picture(str(p), x, y, w, h)
    else:
        box(slide, x, y, w, h, fill=LIGHTBLUE, border=BLUE)
        txt(slide, f"[{p.name}]", x, y + h/2 - Inches(0.3), w, Inches(0.6),
            size=13, color=BLUE, align=PP_ALIGN.CENTER, italic=True)


# ── Slides ──────────────────────────────────────────────────────────────────

def slide_titre(prs):
    s = blank(prs)
    box(s, 0, 0, W, H, fill=NAVY)
    box(s, 0, Inches(5.6), W, Inches(1.9), fill=BLUE)

    # Ligne décorative
    box(s, Inches(0.5), Inches(1.8), Inches(1.2), Inches(0.08), fill=YELLOW)

    txt(s, "Bayesian Sports Analytics",
        Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.5),
        size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    txt(s, "Prédire les résultats de Premier League\navec l'inférence bayésienne",
        Inches(0.5), Inches(3.6), Inches(12.3), Inches(1.2),
        size=24, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.LEFT)

    txt(s, "Cian Higgins  ·  Jules Dantin  ·  Hugo Ferré",
        Inches(0.5), Inches(5.75), Inches(12.3), Inches(0.55),
        size=20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    txt(s, "ECE Paris · IA Probabiliste 2026 · Groupe 02",
        Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.45),
        size=15, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.LEFT)


def slide_probleme(prs):
    s = blank(prs)
    header(s, "La question centrale")
    footer(s)

    # Grande question centrale
    box(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.2), fill=LIGHTBLUE, border=BLUE)
    txt(s, "Peut-on prédire les résultats de football\nmieux que les bookmakers ?",
        Inches(0.7), Inches(1.6), Inches(11.9), Inches(2.0),
        size=30, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # 3 sous-points
    for i, (icon, label, detail) in enumerate([
        ("📊", "Données",   "3 saisons · 1 100 matchs · Premier League"),
        ("🎯", "Modèle",    "Distribution de Poisson + priors bayésiens"),
        ("💰", "Objectif",  "Trouver des paris à valeur positive (value bets)"),
    ]):
        bx = Inches(0.5) + i * Inches(4.28)
        box(s, bx, Inches(4.0), Inches(4.0), Inches(2.9), fill=WHITE, border=BLUE)
        txt(s, icon, bx, Inches(4.1), Inches(4.0), Inches(0.8),
            size=32, align=PP_ALIGN.CENTER)
        txt(s, label, bx, Inches(4.9), Inches(4.0), Inches(0.55),
            size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        txt(s, detail, bx, Inches(5.45), Inches(4.0), Inches(0.7),
            size=13, color=GRAY, align=PP_ALIGN.CENTER)


def slide_donnees(prs):
    s = blank(prs)
    header(s, "Les données", "football-data.co.uk — accès libre, aucune clé API")
    footer(s)

    # 4 stat boxes
    for val, lbl, col, x in [
        ("1 100",  "matchs",           NAVY, Inches(0.5)),
        ("20",     "équipes/saison",   BLUE, Inches(3.6)),
        ("3",      "saisons 2022–25",  BLUE, Inches(6.7)),
        ("5%",     "marge Bet365",     RED,  Inches(9.8)),
    ]:
        stat_box(s, val, lbl, x, Inches(1.5), bg=col)

    # Split train/test
    box(s, Inches(0.5), Inches(3.4), Inches(12.3), Inches(1.7), fill=LIGHTBLUE)
    txt(s, "Split train / test",
        Inches(0.7), Inches(3.45), Inches(12), Inches(0.5),
        size=16, bold=True, color=NAVY)
    txt(s,
        "TRAIN : 3 saisons complètes sauf les 5 dernières matchweeks  →  1 100 matchs",
        Inches(0.7), Inches(3.9), Inches(12), Inches(0.45),
        size=16, color=GRAY)
    txt(s,
        "TEST  : 5 dernières matchweeks 2024-25  →  20 matchs  (simulation prédiction live)",
        Inches(0.7), Inches(4.3), Inches(12), Inches(0.45),
        size=16, color=GRAY)

    # Variables
    bullets(s, [
        "FTHG / FTAG — buts domicile et extérieur (full-time)",
        "FTR — résultat final : H (home), D (draw), A (away)",
        "B365H / B365D / B365A — cotes Bet365 pour chaque issue",
    ], Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.8), base_size=18)


def slide_poisson(prs):
    s = blank(prs)
    header(s, "Pourquoi la loi de Poisson ?",
           "Les buts sont des événements rares et indépendants → Poisson est le bon choix")
    footer(s)

    img(s, FIGURES / "poisson_fit.png",
        Inches(6.8), Inches(1.65), Inches(6.1), Inches(5.2))

    bullets(s, [
        "Un but = événement rare sur 90 minutes",
        "  → Poisson modélise exactement ça",
        "",
        "Paramètre λ = nombre moyen de buts",
        "  → λ_domicile ≈ 1.50  /  λ_extérieur ≈ 1.22",
        "",
        "Test KS : p > 0.05",
        "  → Poisson validé statistiquement",
        "",
        "Buts dom. et ext. peu corrélés",
        "  → On peut les modéliser indépendamment",
    ], Inches(0.5), Inches(1.7), Inches(5.9), Inches(5.3), base_size=19)


def slide_modele(prs):
    s = blank(prs)
    header(s, "Le modèle en 2 formules",
           "Chaque équipe a ses propres paramètres d'attaque et de défense")
    footer(s)

    # Formules dans une boîte
    box(s, Inches(0.5), Inches(1.6), Inches(6.2), Inches(3.8), fill=LIGHTBLUE, border=BLUE)
    txt(s, "Buts domicile  ~  Poisson( θ_dom )\nButs extérieur ~  Poisson( θ_ext )\n\nlog(θ_dom) = μ + δ + attaque[dom] − défense[ext]\nlog(θ_ext) = μ     + attaque[ext] − défense[dom]",
        Inches(0.8), Inches(1.75), Inches(5.7), Inches(3.4),
        size=18, color=NAVY)

    # Légende des paramètres
    for i, (param, explication, color) in enumerate([
        ("μ",          "Base de buts (même pour tous)",     GRAY),
        ("δ",          "Avantage à domicile",                GREEN),
        ("attaque[k]", "Force d'attaque de l'équipe k",     BLUE),
        ("défense[k]", "Solidité défensive de l'équipe k",  RED),
    ]):
        y = Inches(5.65) + i * Inches(0.43)
        box(s, Inches(0.5), y, Inches(1.8), Inches(0.38), fill=LIGHTBLUE)
        txt(s, param, Inches(0.55), y + Pt(3), Inches(1.7), Inches(0.36),
            size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
        txt(s, f"→  {explication}",
            Inches(2.45), y + Pt(3), Inches(9.5), Inches(0.36),
            size=15, color=GRAY)

    # Schéma visuel à droite
    box(s, Inches(7.1), Inches(1.6), Inches(5.7), Inches(3.8), fill=WHITE, border=BLUE)
    txt(s, "Exemple : Arsenal (dom)  vs  Chelsea (ext)",
        Inches(7.2), Inches(1.7), Inches(5.5), Inches(0.5),
        size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    for label, val, y_pos in [
        ("λ_Arsenal  =  exp( μ + δ + attaque_ARS − défense_CHE )", "≈ 1.8 buts", Inches(2.5)),
        ("λ_Chelsea  =  exp( μ     + attaque_CHE − défense_ARS )", "≈ 1.1 buts", Inches(3.1)),
    ]:
        txt(s, label, Inches(7.2), Inches(1.6) + y_pos - Inches(1.5),
            Inches(4.0), Inches(0.5), size=13, color=GRAY)
        txt(s, val, Inches(11.3), Inches(1.6) + y_pos - Inches(1.5),
            Inches(1.3), Inches(0.5), size=15, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)

    txt(s, "→  P(Arsenal gagne) = 52%\n→  P(Nul)            = 22%\n→  P(Chelsea gagne) = 26%",
        Inches(7.2), Inches(3.6), Inches(5.5), Inches(1.7),
        size=18, bold=True, color=NAVY)


def slide_hierarchique(prs):
    s = blank(prs)
    header(s, "Pourquoi 'hiérarchique' ?",
           "Toutes les équipes partagent de l'information entre elles")
    footer(s)

    # Schéma hiérarchique
    # Niveau 1 : hyper-priors
    box(s, Inches(4.5), Inches(1.55), Inches(4.3), Inches(1.0), fill=NAVY)
    txt(s, "μ_α,  σ_α  (Hyperpriors)",
        Inches(4.5), Inches(1.65), Inches(4.3), Inches(0.8),
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Flèches
    for x in [Inches(2.3), Inches(5.0), Inches(7.7), Inches(10.4)]:
        txt(s, "↓", x, Inches(2.65), Inches(0.8), Inches(0.5),
            size=22, color=BLUE, align=PP_ALIGN.CENTER)

    # Niveau 2 : équipes
    team_names = ["Arsenal", "Man City", "Chelsea", "Liverpool"]
    for i, team in enumerate(team_names):
        bx = Inches(0.5) + i * Inches(3.2)
        box(s, bx, Inches(3.25), Inches(2.8), Inches(0.95), fill=BLUE)
        txt(s, f"attaque[{team}]",
            bx, Inches(3.32), Inches(2.8), Inches(0.8),
            size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Explication
    box(s, Inches(0.5), Inches(4.45), Inches(12.3), Inches(2.65), fill=LIGHTBLUE)
    txt(s, "Ce que ça change concrètement :",
        Inches(0.7), Inches(4.55), Inches(12), Inches(0.5),
        size=17, bold=True, color=NAVY)
    bullets(s, [
        "Ipswich vient de monter → peu de matchs dans les données",
        "  → Le modèle s'appuie sur la moyenne de toutes les équipes pour l'estimer",
        "Man City a plein de données → son paramètre est précis",
        "  → Le modèle lui fait davantage confiance",
        "Résultat : pas d'overfitting sur les petites équipes",
    ], Inches(0.7), Inches(5.1), Inches(12), Inches(1.9), base_size=17)


def slide_mcmc(prs):
    s = blank(prs)
    header(s, "Comment on trouve les paramètres ?",
           "MCMC (NUTS) — on explore l'espace des paramètres par simulation")
    footer(s)

    # Schéma simplifié MCMC
    steps = [
        ("1. Proposer", "des valeurs\npour α, β, δ..."),
        ("2. Évaluer", "est-ce que ça\ncorrespond aux data ?"),
        ("3. Accepter\nou rejeter", "selon la probabilité"),
        ("4. Répéter", "8 000 fois\n(2000 × 4 chaînes)"),
    ]
    for i, (titre, detail) in enumerate(steps):
        bx = Inches(0.4) + i * Inches(3.25)
        bg = BLUE if i % 2 == 0 else NAVY
        box(s, bx, Inches(1.55), Inches(2.9), Inches(1.9), fill=bg)
        txt(s, titre, bx, Inches(1.65), Inches(2.9), Inches(0.85),
            size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, detail, bx, Inches(2.5), Inches(2.9), Inches(0.9),
            size=14, color=LIGHTBLUE, align=PP_ALIGN.CENTER)
        if i < 3:
            txt(s, "→", Inches(0.4) + i * Inches(3.25) + Inches(2.95),
                Inches(2.2), Inches(0.3), Inches(0.5),
                size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Résultats de convergence
    box(s, Inches(0.4), Inches(3.7), Inches(12.5), Inches(1.6), fill=LIGHTBLUE)
    txt(s, "Indicateurs de convergence (obligatoires à vérifier)",
        Inches(0.6), Inches(3.78), Inches(12), Inches(0.5),
        size=16, bold=True, color=NAVY)

    for i, (name, val, ok) in enumerate([
        ("R-hat",            "< 1.01",    "✅ Chaînes convergées"),
        ("ESS",              "> 400",     "✅ Samples indépendants"),
        ("Trace plots",      "chenille",  "✅ Bon mélange"),
        ("Divergences NUTS", "≈ 0",       "✅ Exploration complète"),
    ]):
        bx = Inches(0.6) + i * Inches(3.1)
        txt(s, f"{name} {val}", bx, Inches(4.35), Inches(3.0), Inches(0.45),
            size=15, bold=True, color=NAVY)
        txt(s, ok, bx, Inches(4.8), Inches(3.0), Inches(0.4),
            size=14, color=GREEN)

    img(s, FIGURES / "trace_plots.png",
        Inches(0.4), Inches(5.5), Inches(12.5), Inches(1.7))


def slide_home_advantage(prs):
    s = blank(prs)
    header(s, "L'avantage à domicile est réel",
           "Le modèle le quantifie avec son incertitude")
    footer(s)

    img(s, FIGURES / "home_advantage_posterior.png",
        Inches(6.6), Inches(1.5), Inches(6.4), Inches(5.5))

    stat_box(s, "> 99.5%", "P( δ > 0 )", Inches(0.5), Inches(1.6), bg=GREEN)
    stat_box(s, "0.14",    "valeur moyenne de δ", Inches(0.5), Inches(3.45), bg=BLUE)
    stat_box(s, "+15%",    "de buts en plus\nà domicile", Inches(0.5), Inches(5.3), bg=NAVY)

    txt(s, "Interprétation :",
        Inches(3.7), Inches(1.6), Inches(2.7), Inches(0.45),
        size=16, bold=True, color=NAVY)
    txt(s, "exp(0.14) ≈ 1.15\n→ 15% de buts supplémentaires\nquand on joue à domicile",
        Inches(3.7), Inches(2.1), Inches(2.7), Inches(1.5),
        size=16, color=GRAY)

    txt(s, "HDI 94% : [0.08, 0.20]",
        Inches(3.7), Inches(3.8), Inches(2.7), Inches(0.5),
        size=15, italic=True, color=GRAY)
    txt(s, "L'intervalle ne contient pas 0\n→ effet statistiquement robuste",
        Inches(3.7), Inches(4.3), Inches(2.7), Inches(1.0),
        size=15, color=GRAY)


def slide_equipes(prs):
    s = blank(prs)
    header(s, "Force de chaque équipe",
           "Le modèle apprend un paramètre d'attaque ET de défense par équipe")
    footer(s)

    img(s, FIGURES / "attack_defense_quadrant.png",
        Inches(6.8), Inches(1.55), Inches(6.1), Inches(5.4))

    bullets(s, [
        "Attaque positive  →  marque plus que la moyenne",
        "Défense faible    →  encaisse plus que la moyenne",
        "",
        "Top droite : fort en attaque, mauvais en défense",
        "  → Ex : équipes offensives mais fragiles",
        "",
        "Top gauche : fort en attaque ET en défense",
        "  → Les meilleures équipes (Man City, Arsenal...)",
        "",
        "Barres d'erreur = HDI 94% (incertitude du modèle)",
        "  → Équipes avec peu de matchs = plus d'incertitude",
    ], Inches(0.5), Inches(1.65), Inches(5.9), Inches(5.3), base_size=18)


def slide_prediction_match(prs):
    s = blank(prs)
    header(s, "Prédire un match spécifique",
           "Exemple : Man City vs Arsenal — simulation Monte Carlo depuis le posterior")
    footer(s)

    img(s, FIGURES / "scoreline_distribution.png",
        Inches(0.5), Inches(1.55), Inches(6.3), Inches(5.5))

    txt(s, "Comment c'est calculé ?",
        Inches(7.1), Inches(1.65), Inches(5.9), Inches(0.5),
        size=19, bold=True, color=NAVY)

    bullets(s, [
        "On tire 4 000 valeurs depuis le posterior",
        "  → 4 000 valeurs de (attaque, défense, δ...)",
        "",
        "Pour chaque tirage, on simule le match",
        "  → 4 000 scorelines possibles",
        "",
        "On agrège → distribution complète",
        "  → P(2-1), P(1-0), P(0-0)...",
    ], Inches(7.1), Inches(2.3), Inches(5.9), Inches(3.5), base_size=18)

    # Résultat
    box(s, Inches(7.1), Inches(5.95), Inches(5.9), Inches(0.95), fill=LIGHTBLUE, border=BLUE)
    txt(s, "P(Man City)  55%  ·  P(Nul)  22%  ·  P(Arsenal)  23%",
        Inches(7.2), Inches(6.1), Inches(5.7), Inches(0.65),
        size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def slide_evaluation(prs):
    s = blank(prs)
    header(s, "Est-ce que ça marche ?",
           "Évaluation sur les 5 dernières matchweeks de 2024-25 (20 matchs)")
    footer(s)

    # 3 métriques
    for val, lbl, sous, col, x in [
        ("~50%",  "Accuracy",  "Baseline naïf : 43%\nBet365 : 53%",     BLUE,  Inches(0.5)),
        ("~0.20", "RPS",       "Ranked Probability Score\n(plus bas = mieux)", NAVY,  Inches(4.8)),
        ("✓",     "Calibration","Probas prédites ≈\nfréquences observées",    GREEN, Inches(9.1)),
    ]:
        stat_box(s, val, lbl, x, Inches(1.55), bg=col, w=Inches(3.9), h=Inches(1.6))
        txt(s, sous, x, Inches(3.3), Inches(3.9), Inches(0.9),
            size=14, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

    img(s, FIGURES / "calibration_confusion.png",
        Inches(0.5), Inches(4.35), Inches(12.4), Inches(2.7))

    txt(s, "Notre modèle est au niveau des bookmakers sans aucune donnée contextuelle (blessures, météo, etc.)",
        Inches(0.5), Inches(7.05), Inches(12.4), Inches(0.4),
        size=13, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


def slide_value_bets(prs):
    s = blank(prs)
    header(s, "Trouver des paris rentables",
           "Un value bet = notre modèle donne une probabilité plus élevée que le bookmaker")
    footer(s)

    # Schéma conceptuel
    box(s, Inches(0.5), Inches(1.55), Inches(5.8), Inches(2.5), fill=LIGHTBLUE, border=BLUE)
    txt(s, "Bet365 pense : P(Arsenal) = 40%\n→ cote = 1 / 0.40 = 2.50",
        Inches(0.7), Inches(1.65), Inches(5.4), Inches(0.9),
        size=17, color=GRAY)
    txt(s, "Notre modèle : P(Arsenal) = 48%",
        Inches(0.7), Inches(2.5), Inches(5.4), Inches(0.55),
        size=19, bold=True, color=NAVY)
    box(s, Inches(0.7), Inches(3.1), Inches(5.4), Inches(0.75), fill=GREEN)
    txt(s, "EV = 0.48 × 2.50 − 1 = +0.20  →  VALUE BET ✅",
        Inches(0.75), Inches(3.2), Inches(5.3), Inches(0.55),
        size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Kelly
    box(s, Inches(0.5), Inches(4.1), Inches(5.8), Inches(2.95), fill=WHITE, border=NAVY)
    txt(s, "Critère de Kelly",
        Inches(0.7), Inches(4.2), Inches(5.4), Inches(0.5),
        size=18, bold=True, color=NAVY)
    txt(s, "f* = (p × cote − 1) / (cote − 1)\n\nf* = (0.48 × 2.50 − 1) / (2.50 − 1)\n   = 0.20 / 1.50  =  13% du bankroll",
        Inches(0.7), Inches(4.75), Inches(5.4), Inches(2.0),
        size=16, color=GRAY)

    img(s, FIGURES / "model_vs_bookmaker.png",
        Inches(6.7), Inches(1.55), Inches(6.3), Inches(5.5))


def slide_backtest(prs):
    s = blank(prs)
    header(s, "Backtest — simulation sur les données réelles",
           "On rejoue la stratégie sur les 20 matchs de test avec les vrais résultats")
    footer(s)

    img(s, FIGURES / "backtest_results.png",
        Inches(0.5), Inches(1.55), Inches(8.5), Inches(5.5))

    txt(s, "Stratégie Kelly :",
        Inches(9.3), Inches(1.65), Inches(3.7), Inches(0.5),
        size=17, bold=True, color=NAVY)
    bullets(s, [
        "Bankroll initial : 1 000 €",
        "Mise proportionnelle à l'avantage",
        "Cap à 20% du bankroll",
        "",
        "Stratégie mise fixe :",
        "Mise fixe : 20€ par pari",
        "Référence pour comparaison",
        "",
        "Limite importante :",
        "  20 matchs = petit échantillon",
        "  Besoin de plus de données",
        "  pour conclure définitivement",
    ], Inches(9.3), Inches(2.3), Inches(3.7), Inches(4.7), base_size=16)


def slide_dynamique(prs):
    s = blank(prs)
    header(s, "Le modèle dynamique",
           "La force des équipes change au cours de la saison")
    footer(s)

    img(s, FIGURES / "dynamic_attack_trajectories.png",
        Inches(6.6), Inches(1.55), Inches(6.4), Inches(5.5))

    txt(s, "Problème du modèle statique :",
        Inches(0.5), Inches(1.65), Inches(5.8), Inches(0.5),
        size=17, bold=True, color=RED)
    txt(s, "Il suppose qu'Arsenal en janvier\nest aussi fort qu'en août.\nC'est faux !",
        Inches(0.5), Inches(2.2), Inches(5.8), Inches(1.0),
        size=16, color=GRAY)

    txt(s, "Notre solution : Gaussian Random Walk",
        Inches(0.5), Inches(3.4), Inches(5.8), Inches(0.55),
        size=17, bold=True, color=NAVY)

    box(s, Inches(0.5), Inches(4.0), Inches(5.8), Inches(1.0), fill=LIGHTBLUE)
    txt(s, "attaque[équipe, t]  =  attaque[équipe, t-1]  +  bruit",
        Inches(0.65), Inches(4.1), Inches(5.5), Inches(0.8),
        size=15, color=NAVY)

    bullets(s, [
        "La force varie doucement semaine/semaine",
        "Capture : blessures, forme, mercato",
        "Pour prédire le prochain match :",
        "  → On utilise la force du dernier matchweek",
    ], Inches(0.5), Inches(5.2), Inches(5.8), Inches(1.9), base_size=17)


def slide_conclusion(prs):
    s = blank(prs)
    header(s, "En résumé")
    footer(s)

    # 3 grandes boîtes résultats
    for i, (titre, contenu, bg) in enumerate([
        ("Ce qu'on a construit",
         "Modèle Poisson hiérarchique bayésien\nAvantage domicile quantifié\nModèle dynamique (GRW)",
         NAVY),
        ("Ce qu'on a trouvé",
         "P(δ > 0) > 99.5% — home advantage réel\nPrécision ~50% (≈ Bet365)\nValue bets à EV positif identifiés",
         BLUE),
        ("Les limites",
         "Indépendance dom/ext (hypothèse simpliste)\nPas de features contextuels (blessures...)\nTest set petit : 20 matchs seulement",
         ORANGE),
    ]):
        bx = Inches(0.4) + i * Inches(4.35)
        box(s, bx, Inches(1.55), Inches(4.05), Inches(4.0), fill=bg)
        txt(s, titre, bx, Inches(1.65), Inches(4.05), Inches(0.7),
            size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        box(s, bx + Inches(0.1), Inches(2.45), Inches(3.85), Inches(2.9), fill=WHITE)
        txt(s, contenu, bx + Inches(0.2), Inches(2.55), Inches(3.65), Inches(2.7),
            size=15, color=GRAY, wrap=True)

    # Perspectives
    box(s, Inches(0.4), Inches(5.75), Inches(12.5), Inches(1.35), fill=LIGHTBLUE)
    txt(s, "Prochaines étapes :",
        Inches(0.6), Inches(5.85), Inches(3), Inches(0.4),
        size=15, bold=True, color=NAVY)
    txt(s, "Correction Dixon-Coles  ·  Features contextuels (blessures, météo)  ·  Mise à jour en temps réel  ·  Extension multi-championnats",
        Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.6),
        size=14, color=GRAY)


def slide_merci(prs):
    s = blank(prs)
    box(s, 0, 0, W, H, fill=NAVY)
    box(s, 0, Inches(5.5), W, Inches(2.0), fill=BLUE)
    box(s, Inches(0.5), Inches(2.3), Inches(1.5), Inches(0.1), fill=YELLOW)

    txt(s, "Merci !",
        Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5),
        size=64, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    txt(s, "Des questions ?",
        Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.9),
        size=28, color=LIGHTBLUE, align=PP_ALIGN.LEFT)

    txt(s, "Cian Higgins (cian04)  ·  Jules Dantin (julesd13)  ·  Hugo Ferré (hugo-frr)",
        Inches(0.5), Inches(5.65), Inches(12.3), Inches(0.55),
        size=18, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    txt(s, "ECE Paris · IA Probabiliste 2026 · Sujet A.2",
        Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.45),
        size=14, color=LIGHTBLUE, align=PP_ALIGN.LEFT)


# ── Main ─────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()

    slide_titre(prs)
    slide_probleme(prs)
    slide_donnees(prs)
    slide_poisson(prs)
    slide_modele(prs)
    slide_hierarchique(prs)
    slide_mcmc(prs)
    slide_home_advantage(prs)
    slide_equipes(prs)
    slide_prediction_match(prs)
    slide_evaluation(prs)
    slide_value_bets(prs)
    slide_backtest(prs)
    slide_dynamique(prs)
    slide_conclusion(prs)
    slide_merci(prs)

    return prs


if __name__ == "__main__":
    out = Path(__file__).parent / "groupe-02-bayesian-sports-analytics.pptx"
    prs = build()
    prs.save(str(out))
    print(f"✓ {len(prs.slides)} slides générées → {out}")
